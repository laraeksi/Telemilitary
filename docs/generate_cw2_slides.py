from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE_COLOR = RGBColor(15, 23, 42)
ACCENT_COLOR = RGBColor(37, 99, 235)
BODY_COLOR = RGBColor(30, 41, 59)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

    title_tf = slide.shapes.title.text_frame
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = TITLE_COLOR

    sub_tf = slide.placeholders[1].text_frame
    sub_tf.paragraphs[0].font.size = Pt(20)
    sub_tf.paragraphs[0].font.color.rgb = BODY_COLOR


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str], note: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    title_tf = slide.shapes.title.text_frame
    title_tf.paragraphs[0].font.color.rgb = TITLE_COLOR
    title_tf.paragraphs[0].font.bold = True

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for idx, line in enumerate(bullets):
        p = body.add_paragraph() if idx > 0 else body.paragraphs[0]
        p.text = line
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = BODY_COLOR

    if note:
        textbox = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12), Inches(0.6))
        tf = textbox.text_frame
        tf.text = note
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.italic = True
        tf.paragraphs[0].font.color.rgb = ACCENT_COLOR
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    note: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    table_shape = slide.shapes.add_table(
        rows=len(rows) + 1,
        cols=len(headers),
        left=Inches(0.4),
        top=Inches(1.4),
        width=Inches(12.5),
        height=Inches(4.8),
    )
    table = table_shape.table

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = BODY_COLOR

    if note:
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12), Inches(0.7))
        tf = textbox.text_frame
        tf.text = note
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.italic = True
        tf.paragraphs[0].font.color.rgb = ACCENT_COLOR


def add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.9), Inches(1.3), width=Inches(11.5))
    else:
        box = slide.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.2))
        tf = box.text_frame
        tf.text = f"Image missing: {image_path.name}"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(185, 28, 28)

    cap = slide.shapes.add_textbox(Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.5))
    cap_tf = cap.text_frame
    cap_tf.text = caption
    cap_tf.paragraphs[0].font.size = Pt(13)
    cap_tf.paragraphs[0].font.italic = True
    cap_tf.paragraphs[0].font.color.rgb = ACCENT_COLOR


def main() -> None:
    project_root = Path(__file__).resolve().parents[1]
    docs_dir = project_root / "docs"
    output_path = docs_dir / "CW2_Second_Submission_Slides.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Telemilitary - CW2 Final Submission",
        "COMM2020 Project 7 | Indie Game Telemetry and Fair Progression Balancer",
    )

    add_bullets_slide(
        prs,
        "1. What We Delivered in Sprint 2",
        [
            "Playable browser memory game with 10 stages across easy/balanced/hard configs",
            "Telemetry pipeline captures and validates gameplay events in Flask + SQLite",
            "Designer dashboard supports funnel, spikes, progression, fairness, and comparison views",
            "Rule-based balancing toolkit with simulation and decision logging",
        ],
        note="Live demo focus: game -> telemetry -> dashboard -> balancing decision.",
    )

    add_table_slide(
        prs,
        "2. CW2 Specification Alignment (Core)",
        ["Spec requirement", "Project status", "Evidence in this build"],
        [
            ["10+ stages + 3 mechanics", "Met", "10-stage progression, timer/moves, powerups"],
            ["12+ telemetry event types", "Met", "15 documented event types in shared schema"],
            ["Dashboard with fairness + comparisons", "Met", "Dashboard pages and metrics endpoints"],
            ["6+ balancing rules + simulation", "Met", "6 rules in balancing logic + /simulate route"],
            ["RBAC and pseudonymous telemetry", "Met", "Designer login and role checks"],
        ],
        note="All points trace to current repository implementation and docs.",
    )

    add_bullets_slide(
        prs,
        "3. Gameplay and Telemetry Flow",
        [
            "Player sessions emit stage_start, stage_complete/fail, retries, resources, and quit events",
            "Backend validation catches missing fields, impossible sequences, and invalid payload values",
            "Telemetry is stored pseudonymously and reused for seeded analytics even offline",
            "Post-run summary gives players completion/time/resource feedback",
        ],
        note="Talk track: show one short run, then open dashboard immediately.",
    )

    architecture_image = docs_dir / "architecture-diagram.presenter.png"
    if not architecture_image.exists():
        architecture_image = docs_dir / "architecture-diagram.png"

    data_model_image = docs_dir / "data-model-diagram.presenter.png"
    if not data_model_image.exists():
        data_model_image = docs_dir / "data-model-diagram.png"

    add_image_slide(
        prs,
        "4. Architecture Overview",
        architecture_image,
        "Client (React) -> API (Flask) -> SQLite telemetry store -> analytics and balancing endpoints.",
    )

    add_image_slide(
        prs,
        "5. Data Model for Analytics and Auditing",
        data_model_image,
        "Core entities include users, sessions, events, configs, rules, decisions, and anomalies.",
    )

    add_bullets_slide(
        prs,
        "6. Dashboard Insights for Designers",
        [
            "Funnel analytics: stage-by-stage completion and drop-off rates",
            "Difficulty spike detection: high fail/time pressure stages are highlighted",
            "Progression curves: completion time and token economy trends by stage",
            "Fairness indicators: segment comparison (e.g., fast vs slow players)",
            "Config comparison mode for easy vs balanced vs hard",
        ],
    )

    add_bullets_slide(
        prs,
        "7. Balancing Toolkit Demo Scenario",
        [
            "Rule triggers flag problematic stages (e.g., fail rate too high, fairness gap too large)",
            "Designer edits parameters (timer, move limits, helper costs, penalties, rewards)",
            "Simulation estimates before/after completion, failure, and quit rates",
            "Decision log captures change, rationale, and evidence link for auditability",
        ],
        note="Demo example: +10 seconds on Stage 2 increased predicted completion rate.",
    )

    add_table_slide(
        prs,
        "8. Testing and Quality Evidence",
        ["Area", "Automated coverage", "Current evidence"],
        [
            ["Telemetry validation", "Unit tests", "Event required fields and types validated"],
            ["Data ingestion/storage", "Integration tests", "POST /api/events persisted to SQLite"],
            ["Dashboard computations", "Metrics tests", "Funnel/stage/fairness computations verified"],
            ["Balancing rules/simulation", "Balancing tests", "Simulation route returns expected deltas"],
            ["Export endpoints", "API tests", "CSV export route behavior covered"],
        ],
        note="Repository currently contains 15 automated tests across backend test modules.",
    )

    add_bullets_slide(
        prs,
        "9. Client Handover and Maintainability",
        [
            "Deployment guide includes local run, production setup, and test commands",
            "System remains usable with seeded telemetry data if external services are unavailable",
            "Environment-variable based configuration avoids hardcoded secrets",
            "Maintainer extension path: add new event type + update schema, ingestion, and metrics",
        ],
    )

    add_bullets_slide(
        prs,
        "10. Evaluation Highlights and Limitations",
        [
            "Detection validity: rules identify known difficult stages in seeded data",
            "Fairness analysis: segment gaps are measurable and can drive mitigations",
            "Balancing effectiveness: simulation previews expected progression improvements",
            "Limitation: simulation is lightweight and should be validated with further playtesting",
            "Limitation: final manual testing evidence should be packaged in testing_evidence.pdf",
        ],
    )

    add_bullets_slide(
        prs,
        "11. Risks, Next Steps, and Closing",
        [
            "Risk: over-tuning for one segment can hurt fairness for others",
            "Risk: telemetry quality degrades insight quality if validation rules are bypassed",
            "Next step: continue scenario-based balancing cycles with decision audit trail",
            "Next step: extend manual E2E evidence pack and maintainer runbook details",
            "Final outcome: Telemilitary provides a practical telemetry-driven balancing workflow",
        ],
        note="Final minute in presentation: recap value, limitations, and practical roadmap.",
    )

    prs.save(output_path)
    print(f"Created: {output_path}")


if __name__ == "__main__":
    main()
